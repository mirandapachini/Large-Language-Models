"""Core document ingestion, retrieval, prompting, and Ollama utilities.

The module intentionally avoids LangChain so that students can see each RAG step:
1. read course documents,
2. split them into overlapping chunks,
3. create embeddings,
4. retrieve semantically similar chunks,
5. send retrieved evidence to a local LLM.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence
import hashlib
import re

import numpy as np
import pandas as pd
import pymupdf
import requests
from docx import Document
from pptx import Presentation

try:
    import faiss  # type: ignore
except ImportError:  # The app has a NumPy fallback for easier Windows installation.
    faiss = None


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv", ".xlsx"}
DEFAULT_EMBEDDING_MODEL = "all-minilm:latest"


@dataclass(frozen=True)
class TextChunk:
    """A retrievable piece of a source document."""

    text: str
    source: str
    locator: str
    chunk_number: int

    @property
    def citation(self) -> str:
        return f"{self.source} | {self.locator}"


@dataclass(frozen=True)
class SearchResult:
    """A retrieved chunk and its cosine-similarity score."""

    chunk: TextChunk
    score: float


@dataclass
class KnowledgeBase:
    """Embedding configuration, vector index, and source chunks."""

    embedding_model_name: str
    ollama_url: str
    chunks: list[TextChunk]
    embeddings: np.ndarray
    faiss_index: object | None = None

    @property
    def backend_name(self) -> str:
        return "FAISS" if self.faiss_index is not None else "NumPy cosine search"

    def search(self, query: str, k: int = 5) -> list[SearchResult]:
        if not self.chunks:
            return []

        requested_k = max(1, min(k, len(self.chunks)))
        query_vector = embed_texts_ollama(
            [query],
            model_name=self.embedding_model_name,
            base_url=self.ollama_url,
        )

        if self.faiss_index is not None:
            scores, indices = self.faiss_index.search(query_vector, requested_k)
            ranked = zip(scores[0].tolist(), indices[0].tolist())
        else:
            # Embeddings are normalized, so the dot product equals cosine similarity.
            similarities = self.embeddings @ query_vector[0]
            indices = np.argsort(similarities)[::-1][:requested_k]
            ranked = ((float(similarities[i]), int(i)) for i in indices)

        results: list[SearchResult] = []
        for score, index in ranked:
            if index < 0:
                continue
            results.append(SearchResult(chunk=self.chunks[index], score=float(score)))
        return results


def clean_text(text: str) -> str:
    """Normalize whitespace while preserving readable paragraph breaks."""

    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n[ \t]+", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, chunk_size: int = 180, overlap: int = 35) -> list[str]:
    """Split text into overlapping word chunks.

    all-MiniLM-L6-v2 has a relatively short token limit. Compact chunks reduce
    truncation and usually improve retrieval precision for course notes.
    """

    if chunk_size <= 0:
        raise ValueError("chunk_size must be positive")
    if overlap < 0 or overlap >= chunk_size:
        raise ValueError("overlap must be non-negative and smaller than chunk_size")

    words = clean_text(text).split()
    if not words:
        return []

    step = chunk_size - overlap
    chunks: list[str] = []
    for start in range(0, len(words), step):
        piece = words[start : start + chunk_size]
        if not piece:
            break
        chunks.append(" ".join(piece))
        if start + chunk_size >= len(words):
            break
    return chunks


def _make_chunks(
    text: str,
    source: str,
    locator: str,
    chunk_size: int,
    overlap: int,
) -> list[TextChunk]:
    pieces = chunk_text(text, chunk_size=chunk_size, overlap=overlap)
    return [
        TextChunk(text=piece, source=source, locator=locator, chunk_number=i + 1)
        for i, piece in enumerate(pieces)
    ]


def _read_pdf(path: Path, chunk_size: int, overlap: int) -> list[TextChunk]:
    chunks: list[TextChunk] = []
    with pymupdf.open(path) as document:
        for page_number, page in enumerate(document, start=1):
            text = page.get_text("text")
            chunks.extend(
                _make_chunks(
                    text,
                    source=path.name,
                    locator=f"page {page_number}",
                    chunk_size=chunk_size,
                    overlap=overlap,
                )
            )
    return chunks


def _read_docx(path: Path, chunk_size: int, overlap: int) -> list[TextChunk]:
    document = Document(path)
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]

    for table_number, table in enumerate(document.tables, start=1):
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        if rows:
            parts.append(f"Table {table_number}\n" + "\n".join(rows))

    return _make_chunks(
        "\n\n".join(parts),
        source=path.name,
        locator="document",
        chunk_size=chunk_size,
        overlap=overlap,
    )


def _read_pptx(path: Path, chunk_size: int, overlap: int) -> list[TextChunk]:
    presentation = Presentation(path)
    chunks: list[TextChunk] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                table_rows = [
                    " | ".join(cell.text.strip() for cell in row.cells)
                    for row in shape.table.rows
                ]
                texts.extend(table_rows)
        chunks.extend(
            _make_chunks(
                "\n".join(texts),
                source=path.name,
                locator=f"slide {slide_number}",
                chunk_size=chunk_size,
                overlap=overlap,
            )
        )
    return chunks


def _read_text(path: Path, chunk_size: int, overlap: int) -> list[TextChunk]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return _make_chunks(
        text,
        source=path.name,
        locator="document",
        chunk_size=chunk_size,
        overlap=overlap,
    )


def _dataframe_chunks(
    frame: pd.DataFrame,
    source: str,
    locator_prefix: str,
    rows_per_chunk: int = 20,
) -> list[TextChunk]:
    chunks: list[TextChunk] = []
    if frame.empty:
        return chunks

    # RAG can retrieve table excerpts, but it is not a substitute for a full
    # statistical computation engine. Limit very large files to keep indexing sane.
    frame = frame.head(5000)
    for chunk_number, start in enumerate(range(0, len(frame), rows_per_chunk), start=1):
        subset = frame.iloc[start : start + rows_per_chunk]
        text = subset.to_csv(index=False)
        end = start + len(subset)
        chunks.append(
            TextChunk(
                text=clean_text(text),
                source=source,
                locator=f"{locator_prefix}, rows {start + 1}-{end}",
                chunk_number=chunk_number,
            )
        )
    return chunks


def _read_csv(path: Path, _chunk_size: int, _overlap: int) -> list[TextChunk]:
    frame = pd.read_csv(path)
    return _dataframe_chunks(frame, path.name, "table")


def _read_xlsx(path: Path, _chunk_size: int, _overlap: int) -> list[TextChunk]:
    sheets = pd.read_excel(path, sheet_name=None)
    chunks: list[TextChunk] = []
    for sheet_name, frame in sheets.items():
        chunks.extend(_dataframe_chunks(frame, path.name, f"sheet {sheet_name}"))
    return chunks


READERS = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".txt": _read_text,
    ".md": _read_text,
    ".csv": _read_csv,
    ".xlsx": _read_xlsx,
}


def discover_documents(docs_dir: Path) -> list[Path]:
    """Return supported documents in deterministic filename order."""

    if not docs_dir.exists():
        return []
    return sorted(
        [
            path
            for path in docs_dir.iterdir()
            if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
        ],
        key=lambda p: p.name.lower(),
    )


def document_signature(docs_dir: Path) -> str:
    """Create a stable cache key based on filenames, sizes, and modification times."""

    digest = hashlib.sha256()
    for path in discover_documents(docs_dir):
        stat = path.stat()
        digest.update(path.name.encode("utf-8", errors="ignore"))
        digest.update(str(stat.st_size).encode())
        digest.update(str(stat.st_mtime_ns).encode())
    return digest.hexdigest()


def load_document_chunks(
    docs_dir: Path,
    chunk_size: int = 180,
    overlap: int = 35,
) -> tuple[list[TextChunk], list[str]]:
    """Parse all supported course documents and report nonfatal file errors."""

    chunks: list[TextChunk] = []
    errors: list[str] = []
    for path in discover_documents(docs_dir):
        reader = READERS[path.suffix.lower()]
        try:
            chunks.extend(reader(path, chunk_size, overlap))
        except Exception as exc:  # Continue indexing other files.
            errors.append(f"{path.name}: {type(exc).__name__}: {exc}")
    return chunks, errors


def embed_texts_ollama(
    texts: Sequence[str],
    model_name: str = DEFAULT_EMBEDDING_MODEL,
    base_url: str = "http://localhost:11434",
    batch_size: int = 64,
    timeout: int = 240,
) -> np.ndarray:
    """Generate normalized embeddings through Ollama's local /api/embed endpoint."""

    if not texts:
        return np.empty((0, 0), dtype="float32")

    vectors: list[list[float]] = []
    for start in range(0, len(texts), batch_size):
        batch = list(texts[start : start + batch_size])
        try:
            response = requests.post(
                f"{base_url.rstrip('/')}/api/embed",
                json={"model": model_name, "input": batch, "truncate": True},
                timeout=timeout,
            )
            response.raise_for_status()
            data = response.json()
        except requests.Timeout as exc:
            raise RuntimeError(
                "Embedding generation timed out. Try fewer or smaller course files."
            ) from exc
        except requests.RequestException as exc:
            detail = ""
            if getattr(exc, "response", None) is not None:
                detail = f" Server response: {exc.response.text[:500]}"
            raise RuntimeError(
                f"Ollama embedding request failed for model '{model_name}': {exc}.{detail}"
            ) from exc

        batch_vectors = data.get("embeddings")
        if not batch_vectors or len(batch_vectors) != len(batch):
            raise RuntimeError(
                f"Ollama returned an invalid embedding response for '{model_name}'."
            )
        vectors.extend(batch_vectors)

    matrix = np.asarray(vectors, dtype="float32")
    norms = np.linalg.norm(matrix, axis=1, keepdims=True)
    matrix = matrix / np.maximum(norms, 1e-12)
    return matrix


def build_knowledge_base(
    docs_dir: Path,
    embedding_model_name: str = DEFAULT_EMBEDDING_MODEL,
    ollama_url: str = "http://localhost:11434",
    chunk_size: int = 180,
    overlap: int = 35,
    prefer_faiss: bool = True,
) -> tuple[KnowledgeBase, list[str]]:
    """Read documents, embed chunks, and construct a cosine-similarity index."""

    chunks, errors = load_document_chunks(docs_dir, chunk_size, overlap)
    if not chunks:
        raise ValueError(
            "No readable text was found. Add a supported course document to the docs folder."
        )

    texts = [chunk.text for chunk in chunks]
    embeddings = embed_texts_ollama(
        texts,
        model_name=embedding_model_name,
        base_url=ollama_url,
    )

    index = None
    if prefer_faiss and faiss is not None:
        index = faiss.IndexFlatIP(embeddings.shape[1])
        index.add(embeddings)

    return (
        KnowledgeBase(
            embedding_model_name=embedding_model_name,
            ollama_url=ollama_url,
            chunks=chunks,
            embeddings=embeddings,
            faiss_index=index,
        ),
        errors,
    )


def format_retrieved_context(results: Sequence[SearchResult]) -> str:
    """Format evidence with source labels the model can cite exactly."""

    sections: list[str] = []
    for number, result in enumerate(results, start=1):
        sections.append(
            f"[SOURCE {number}: {result.chunk.citation}]\n{result.chunk.text}"
        )
    return "\n\n".join(sections)


def build_tutor_messages(
    question: str,
    results: Sequence[SearchResult],
    response_style: str = "Teach step by step",
    recent_history: Iterable[dict[str, str]] | None = None,
) -> list[dict[str, str]]:
    """Create a grounded chat prompt for an OnboardBot."""

    system_message = """You are OnboardBot, a friendly and knowledgeable HR onboarding assistant for new employees.

Follow these rules:
1. Base company-specific claims on the supplied retrieved context. Never invent a policy, document, deadline, dollar amount, or quotation.
2. Cite supporting material using the exact labels supplied, for example [SOURCE 2].
3. When the context is insufficient, state that clearly. You may then note that the employee should confirm with HR or IT, but do not pretend an answer came from the company documents if it did not.
4. Explain the policy or process first in plain language, then give the specific steps or numbers, then note any relevant deadline or next action the employee should take.
5. For anything involving dates, dollar amounts, or percentages, quote the exact figure from the retrieved context and double-check it before answering.
6. Distinguish between company-wide policy and anything that depends on individual circumstances (e.g., manager approval, role type, tenure), and say when the employee should check with their manager or HR Business Partner.
7. Do not fabricate deadlines, contact information, or numerical values. Ask for missing details when a question cannot be fully answered from the retrieved context.
8. Use readable Markdown, a warm and approachable tone, and keep the language clear for someone in their first days at the company.
"""

    history_text = ""
    if recent_history:
        concise_history = list(recent_history)[-6:]
        history_text = "\n\nRecent conversation:\n" + "\n".join(
            f"{item.get('role', 'user').upper()}: {item.get('content', '')}"
            for item in concise_history
        )

    context = format_retrieved_context(results)
    user_message = f"""Response style: {response_style}

Retrieved course context:
{context or '[No relevant course context was retrieved.]'}
{history_text}

Student question:
{question}

Answer the student using the retrieved evidence and the tutoring rules."""

    return [
        {"role": "system", "content": system_message},
        {"role": "user", "content": user_message},
    ]


def check_ollama(base_url: str = "http://localhost:11434") -> tuple[bool, list[str], str]:
    """Check whether Ollama is reachable and return installed model names."""

    try:
        response = requests.get(f"{base_url.rstrip('/')}/api/tags", timeout=4)
        response.raise_for_status()
        payload = response.json()
        models = [item.get("name", "") for item in payload.get("models", []) if item.get("name")]
        return True, models, "Ollama is running."
    except requests.RequestException as exc:
        return False, [], f"Ollama is not reachable: {exc}"


def query_ollama(
    messages: Sequence[dict[str, str]],
    model_name: str = "llama3.2:latest",
    base_url: str = "http://localhost:11434",
    temperature: float = 0.2,
    timeout: int = 240,
) -> str:
    """Send a non-streaming chat request to the local Ollama server."""

    payload = {
        "model": model_name,
        "messages": list(messages),
        "stream": False,
        "options": {
            "temperature": temperature,
            "num_ctx": 4096,
        },
    }
    try:
        response = requests.post(
            f"{base_url.rstrip('/')}/api/chat",
            json=payload,
            timeout=timeout,
        )
        response.raise_for_status()
        data = response.json()
    except requests.Timeout as exc:
        raise RuntimeError(
            "The model request timed out. Try a shorter question or close other memory-intensive programs."
        ) from exc
    except requests.RequestException as exc:
        detail = ""
        if getattr(exc, "response", None) is not None:
            detail = f" Server response: {exc.response.text[:500]}"
        raise RuntimeError(f"Ollama request failed: {exc}.{detail}") from exc

    answer = data.get("message", {}).get("content", "").strip()
    if not answer:
        raise RuntimeError("Ollama returned an empty response.")
    return answer
